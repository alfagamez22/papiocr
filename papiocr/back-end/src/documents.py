"""Document translation: extract text, translate, rebuild.

Handles DOCX, PDF, XLSX, PPTX. Auto-detects scanned vs digital:
tries native text extraction first; falls back to OCR pipeline if empty.
"""
from __future__ import annotations

import io
import os
import tempfile
from pathlib import Path
from typing import TYPE_CHECKING

from .config import OUTPUT_DIR
from .utils import LOGGER, ensure_dir

if TYPE_CHECKING:
    from .pipeline import ImageTranslationPipeline

SUPPORTED_DOC_FORMATS = {".docx", ".pdf", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"}
MAX_SIZE_MB = 15


# ---------------------------------------------------------------------------
# Native text extraction (digital docs)
# ---------------------------------------------------------------------------


def _extract_text_docx(path: Path) -> list[str]:
    import docx  # type: ignore
    doc = docx.Document(str(path))
    return [p.text for p in doc.paragraphs if p.text.strip()]


def _extract_text_pdf(path: Path) -> list[str]:
    import fitz  # type: ignore
    doc = fitz.open(str(path))
    texts = []
    for page in doc:
        t = page.get_text().strip()
        if t:
            texts.append(t)
    return texts


def _extract_text_xlsx(path: Path) -> list[tuple[str, str, str]]:
    """Return (sheet, cell_ref, text) for every non-empty cell."""
    import openpyxl  # type: ignore
    wb = openpyxl.load_workbook(path, data_only=True)
    entries: list[tuple[str, str, str]] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    entries.append((sheet_name, cell.coordinate, cell.value.strip()))
    return entries


def _extract_text_pptx(path: Path) -> list[tuple[int, str]]:
    """Return (slide_num, text) for every text-bearing shape."""
    from pptx import Presentation  # type: ignore
    prs = Presentation(str(path))
    entries: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        entries.append((i, t))
    return entries


# ---------------------------------------------------------------------------
# Rebuild documents with translated text
# ---------------------------------------------------------------------------


def _rebuild_docx(original: Path, translations: dict[str, str]) -> Path:
    import docx
    doc = docx.Document(str(original))
    for para in doc.paragraphs:
        if para.text.strip() in translations:
            para.text = translations[para.text.strip()]
    out = OUTPUT_DIR / f"{original.stem}_translated.docx"
    ensure_dir(out.parent)
    doc.save(str(out))
    return out


def _rebuild_pdf_text(original: Path, translations: dict[str, str], target: str) -> Path:
    """Rebuild a PDF with translated text (digital PDFs only)."""
    import fitz
    doc = fitz.open(str(original))
    for page in doc:
        text_blocks = page.get_text("blocks")
        for block in text_blocks:
            t = block[4].strip()
            if t in translations:
                # Remove original and insert translated in same position
                rect = fitz.Rect(block[:4])
                page.add_redact_annot(rect, fill=(255, 255, 255))
                page.apply_redactions()
                page.insert_text(
                    (rect.x0, rect.y0 + rect.height - 2),
                    translations[t],
                    fontsize=block[3] - block[1] if (block[3] - block[1]) > 0 else 11,
                    color=(0, 0, 0),
                )
    out = OUTPUT_DIR / f"{original.stem}_{target}.pdf"
    ensure_dir(out.parent)
    doc.save(str(out))
    return out


def _rebuild_xlsx(original: Path, translations: dict[str, str]) -> Path:
    import openpyxl
    wb = openpyxl.load_workbook(original, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip() in translations:
                    cell.value = translations[cell.value.strip()]
    out = OUTPUT_DIR / f"{original.stem}_translated.xlsx"
    ensure_dir(out.parent)
    wb.save(str(out))
    return out


def _rebuild_pptx(original: Path, translations: dict[str, str]) -> Path:
    from pptx import Presentation
    prs = Presentation(str(original))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t in translations:
                        # Clear existing runs and replace with translated
                        for run in para.runs:
                            run.text = ""
                        if para.runs:
                            para.runs[0].text = translations[t]
                        else:
                            from pptx.util import Pt
                            from pptx.enum.text import PP_ALIGN
                            run = para.add_run()
                            run.text = translations[t]
    out = OUTPUT_DIR / f"{original.stem}_translated.pptx"
    ensure_dir(out.parent)
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# OCR-based document translation (scanned docs / image-rendered pages)
# ---------------------------------------------------------------------------


def _ocr_translate_pdf(
    path: Path, pipe: ImageTranslationPipeline, source: str, target: str
) -> Path:
    """Render PDF pages to images, OCR translate, embed into new PDF."""
    import fitz
    from PIL import Image
    import numpy as np
    import cv2

    pdf_doc = fitz.open(str(path))
    out_doc = fitz.open()
    img_dir = ensure_dir(OUTPUT_DIR / "_pages")

    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        # Render page to image at 200 DPI
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat)
        img_path = img_dir / f"page_{page_num:04d}.png"
        pix.save(str(img_path))

        result = pipe.run(img_path, output_path=img_dir / f"page_{page_num:04d}_translated.png")
        translated_img = result.get("output_path", str(img_dir / f"page_{page_num:04d}_translated.png"))

        # Embed translated image as a new PDF page
        pil_img = Image.open(translated_img)
        img_arr = np.array(pil_img)
        h, w = img_arr.shape[:2]
        rect = fitz.Rect(0, 0, w, h)
        new_page = out_doc.new_page(width=w, height=h)
        new_page.insert_image(rect, filename=translated_img)

    out = OUTPUT_DIR / f"{path.stem}_{target}.pdf"
    out_doc.save(str(out))
    pdf_doc.close()
    out_doc.close()
    return out


# ---------------------------------------------------------------------------
# Auto-detect: native text extraction vs OCR pipeline
# ---------------------------------------------------------------------------


def _should_ocr(text_blocks: list[str]) -> bool:
    combined = "".join(text_blocks).strip()
    if len(combined) < 20:
        return True
    # If most of the "text" is whitespace/non-alphanumeric, assume scanned
    alpha = sum(c.isalpha() or c.isdigit() for c in combined)
    return (alpha / max(len(combined), 1)) < 0.3


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def auto_translate_document(
    path: Path,
    pipe: ImageTranslationPipeline,
    ext: str,
    source: str,
    target: str,
) -> Path:
    if ext in {".doc", ".docx"}:
        texts = _extract_text_docx(path)
        if _should_ocr(texts):
            LOGGER.info("DOCX appears scanned; falling back to image pipeline")
            return _ocr_translate_pdf(path, pipe, source, target)
        translated = {t: pipe.translator.translate(t) for t in texts if t.strip()}
        return _rebuild_docx(path, translated)

    if ext == ".pdf":
        texts = _extract_text_pdf(path)
        if _should_ocr(texts):
            LOGGER.info("PDF appears scanned; falling back to image pipeline")
            return _ocr_translate_pdf(path, pipe, source, target)
        translated = {t: pipe.translator.translate(t) for t in texts if t.strip()}
        return _rebuild_pdf_text(path, translated, target)

    if ext in {".xls", ".xlsx"}:
        entries = _extract_text_xlsx(path)
        texts = [e[2] for e in entries]
        if _should_ocr(texts):
            LOGGER.warning("XLSX has minimal text; continuing with native extraction")
        translated = {t: pipe.translator.translate(t) for t in texts if t.strip()}
        return _rebuild_xlsx(path, translated)

    if ext in {".ppt", ".pptx"}:
        entries = _extract_text_pptx(path)
        texts = [e[1] for e in entries]
        if _should_ocr(texts):
            LOGGER.info("PPTX appears scanned; falling back to image pipeline")
            return _ocr_translate_pdf(path, pipe, source, target)
        translated = {t: pipe.translator.translate(t) for t in texts if t.strip()}
        return _rebuild_pptx(path, translated)

    raise ValueError(f"Unsupported format: {ext}")
